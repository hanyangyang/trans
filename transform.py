# encoding: utf-8

from pptx import Presentation
from googletrans import Translator
import sys 
import os 
import time

import tkinter.messagebox
import tkinter as tk
 
window=tk.Tk()
window.title('menu')
window.geometry('200x200')
window.resizable(False, False)#固定窗体
window.withdraw()
 

translator = Translator(
    service_urls=['translate.google.cn'],
    user_agent='Mozilla/5.0 (Macintosh; Intel Mac OS X 10_14_0) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/70.0.3538.77 Safari/537.36',
    proxies={'http':'http://189.91.231.43:3128'}
    )

text_runs = []
all_runs = []
count = 0

# 此函数检查是不是汉语
def check_contain_chinese(check_str):
    for ch in check_str:
        if u'\u4e00' <= ch <= u'\u9fff':
            return False
    return True

all_items = os.listdir(os.path.dirname(os.path.realpath(sys.argv[0])))
for item in all_items:
    print(item)
    (name,type) = os.path.splitext(os.path.dirname(os.path.realpath(sys.argv[0]))+item)
    if not type==".pptx":
        continue
    print(item)    
    prs = Presentation(os.path.dirname(os.path.realpath(sys.argv[0]))+"/"+item)
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if  check_contain_chinese(run.text):
                        continue
                    text_runs.append(run.text)
                    all_runs.append(run)
                    if len(text_runs)<41:
                        continue
                    count = count + 1
                    print(count)
                    text = translator.translate(text_runs,src='zh-cn',dest='en')
                    for runDetail  in all_runs:
                        index = all_runs.index(runDetail)
                        runDetail.text = runDetail.text+'\r'+text[index].text
                        print(runDetail.text)
                    text_runs = []
                    all_runs = []
                    time.sleep(2)
    text = translator.translate(text_runs,src='zh-cn',dest='en')
    for runDetail  in all_runs:
        index = all_runs.index(runDetail)
        runDetail.text = runDetail.text+'\r'+text[index].text
        print(runDetail.text)
    prs.save(os.path.dirname(os.path.realpath(sys.argv[0]))+"/"+"翻译后"+item)
    time.sleep(10)
tkinter.messagebox.showinfo('来自作者的友情提示', '如果觉得软件不错，您可以选择请肘子吃饭！❤️')

